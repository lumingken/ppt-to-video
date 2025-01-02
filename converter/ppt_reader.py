from typing import List, Tuple
from pptx import Presentation
from PIL import Image
from pathlib import Path
from dataclasses import dataclass
import comtypes.client

@dataclass
class SlideContent:
    image_path: Path
    text: str = ""

class PPTReader:
    def __init__(self, ppt_path):
        self.ppt_path = Path(ppt_path).resolve()  # Get absolute path
        self.presentation = Presentation(ppt_path)
    
    def extract_content(self):
        """
        Extracts content (images and text) from each slide
        Returns a list of SlideContent objects
        """
        slides_content = []
        temp_dir = Path('temp').resolve()  # Get absolute path
        temp_dir.mkdir(exist_ok=True)
        
        # Convert slides to images using PowerPoint
        self._convert_slides_to_images(temp_dir)
        
        # Extract text and create SlideContent objects
        for i, slide in enumerate(self.presentation.slides):
            image_path = temp_dir / f'slide_{i}.png'
            text = self._extract_slide_text(slide)
            slides_content.append(SlideContent(image_path, text))
        
        return slides_content
    
    def _convert_slides_to_images(self, output_dir):
        """Convert PowerPoint slides to images using PowerPoint automation"""
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = True  # This helps avoid some automation issues
        
        try:
            # Open the presentation
            pres = powerpoint.Presentations.Open(str(self.ppt_path))
            
            # Export each slide as PNG
            for i in range(1, pres.Slides.Count + 1):
                slide_path = (output_dir / f'slide_{i-1}.png').resolve()
                # Create the output directory if it doesn't exist
                slide_path.parent.mkdir(parents=True, exist_ok=True)
                # Export the slide
                pres.Slides(i).Export(str(slide_path), "PNG")
            
            # Close the presentation
            pres.Close()
            
        finally:
            # Make sure PowerPoint is closed
            try:
                powerpoint.Quit()
            except:
                pass
    
    def _extract_slide_text(self, slide):
        """Extract text from slide notes"""
        if slide.has_notes_slide and slide.notes_slide:
            # Get text from the notes
            notes_text = slide.notes_slide.notes_text_frame.text
            return notes_text.strip()
        return ""  # Return empty string if no notes 